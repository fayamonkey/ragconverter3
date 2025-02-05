import mimetypes
import pdfplumber
import markdown
from pathlib import Path
from typing import Dict, Any, List, Union
import logging
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from io import BytesIO
from langdetect import detect
import tempfile
import os
import easyocr
import numpy as np
from PIL import Image

class DocumentProcessor:
    """Handles the processing of different document types."""
    
    def __init__(self):
        mimetypes.init()
        self.supported_types = {
            'text/plain': self._process_text,
            'text/markdown': self._process_markdown,
            'application/pdf': self._process_pdf,
            'text/html': self._process_html,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx
        }
        
        # Initialize OCR
        self.ocr_reader = None
        self.ocr_available = False

    def _ensure_ocr_initialized(self):
        """Initialize OCR reader if not already initialized."""
        if self.ocr_reader is None:
            try:
                self.ocr_reader = easyocr.Reader(['en'])
                self.ocr_available = True
            except Exception as e:
                logging.warning(f"OCR initialization failed: {str(e)}")
                self.ocr_available = False

    def _detect_language(self, text: str) -> str:
        """Detect the language of the text."""
        try:
            return detect(text)
        except:
            return 'unknown'

    def _extract_text_from_image(self, image) -> str:
        """Extract text from an image using OCR."""
        if not self.ocr_available:
            self._ensure_ocr_initialized()
            if not self.ocr_available:
                return ""

        try:
            # Convert PIL Image to numpy array if necessary
            if isinstance(image, Image.Image):
                image = np.array(image)

            results = self.ocr_reader.readtext(image)
            return '\n'.join(text for _, text, conf in results if conf > 0.5)
        except Exception as e:
            logging.error(f"OCR error: {str(e)}")
            return ""

    def _process_docx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Word documents."""
        doc = Document(BytesIO(content))
        
        # Extract text from paragraphs and tables
        text_content = []
        for paragraph in doc.paragraphs:
            text_content.append(paragraph.text)
        
        for table in doc.tables:
            for row in table.rows:
                text_content.append(' | '.join(cell.text for cell in row.cells))
        
        full_text = '\n'.join(text_content)
        
        return {
            'content': full_text,
            'metadata': {
                'filename': filename,
                'type': 'docx',
                'size': len(content),
                'language': self._detect_language(full_text)
            }
        }

    def _process_xlsx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Excel spreadsheets."""
        wb = load_workbook(BytesIO(content))
        
        text_content = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_content.append(f"\n### Worksheet: {sheet} ###\n")
            
            for row in ws.iter_rows():
                row_content = []
                for cell in row:
                    if cell.value is not None:
                        row_content.append(str(cell.value))
                if row_content:
                    text_content.append(' | '.join(row_content))
        
        full_text = '\n'.join(text_content)
        
        return {
            'content': full_text,
            'metadata': {
                'filename': filename,
                'type': 'xlsx',
                'size': len(content),
                'sheets': wb.sheetnames,
                'language': self._detect_language(full_text)
            }
        }

    def _process_pptx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint presentations."""
        prs = Presentation(BytesIO(content))
        
        text_content = []
        for slide_number, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n### Slide {slide_number} ###\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
        
        full_text = '\n'.join(text_content)
        
        return {
            'content': full_text,
            'metadata': {
                'filename': filename,
                'type': 'pptx',
                'size': len(content),
                'slides': len(prs.slides),
                'language': self._detect_language(full_text)
            }
        }

    def _process_pdf(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF files with OCR support."""
        text_content = []
        metadata = {}
        images_processed = False
        
        # First try normal text extraction
        with BytesIO(content) as pdf_file:
            with pdfplumber.open(pdf_file) as pdf:
                metadata = pdf.metadata
                for page in pdf.pages:
                    page_text = page.extract_text() or ''
                    text_content.append(page_text)
                    
                    # If page has little or no text, try OCR
                    if len(page_text.strip()) < 100:
                        try:
                            # Convert page to image
                            img = page.to_image()
                            img_array = np.array(img.original)
                            
                            # Perform OCR
                            ocr_text = self._extract_text_from_image(img_array)
                            if ocr_text.strip():
                                text_content.append(ocr_text)
                                images_processed = True
                        except Exception as e:
                            logging.error(f"Error during OCR processing: {str(e)}")
        
        full_text = '\n'.join(text_content)
        
        return {
            'content': full_text,
            'metadata': {
                'filename': filename,
                'type': 'pdf',
                'size': len(content),
                'pages': len(text_content),
                'pdf_metadata': metadata,
                'ocr_applied': images_processed,
                'language': self._detect_language(full_text)
            }
        }
    
    def _process_text(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process plain text files."""
        text = content.decode('utf-8', errors='ignore')
        return {
            'content': text,
            'metadata': {
                'filename': filename,
                'type': 'text',
                'size': len(content)
            }
        }
    
    def _process_markdown(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process markdown files."""
        text = content.decode('utf-8', errors='ignore')
        html = markdown.markdown(text)
        return {
            'content': text,
            'html_content': html,
            'metadata': {
                'filename': filename,
                'type': 'markdown',
                'size': len(content)
            }
        }
    
    def _process_html(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process HTML files."""
        text = content.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(text, 'html.parser')
        
        # Extract text content
        text_content = soup.get_text(separator='\n', strip=True)
        
        # Extract metadata from meta tags
        meta_tags = {}
        for meta in soup.find_all('meta'):
            name = meta.get('name', meta.get('property', ''))
            content = meta.get('content', '')
            if name and content:
                meta_tags[name] = content
        
        return {
            'content': text_content,
            'metadata': {
                'filename': filename,
                'type': 'html',
                'size': len(content),
                'meta_tags': meta_tags
            }
        }

    def process_document(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Process a document and return its content and metadata.
        
        Args:
            file_content: The binary content of the file
            filename: Name of the file
            
        Returns:
            Dict containing processed content and metadata
        """
        try:
            # Determine file type based on extension
            mime_type, _ = mimetypes.guess_type(filename)
            
            # Handle specific file types
            if filename.lower().endswith('.md'):
                mime_type = 'text/markdown'
            elif filename.lower().endswith('.docx'):
                mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif filename.lower().endswith('.xlsx'):
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif filename.lower().endswith('.pptx'):
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            # Default to text/plain for unknown types
            elif not mime_type:
                mime_type = 'text/plain'
            
            if mime_type not in self.supported_types:
                raise ValueError(f"Unsupported file type: {mime_type}")
            
            processor = self.supported_types[mime_type]
            return processor(file_content, filename)
            
        except Exception as e:
            logging.error(f"Error processing document {filename}: {str(e)}")
            raise 